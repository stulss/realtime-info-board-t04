from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

ROOT=Path(__file__).resolve().parents[1]
OUT=ROOT/'docs'/'T04_제출자료_v2_2026-08-27.pptx'
if OUT.exists(): raise SystemExit(f'refusing to overwrite {OUT}')
S=ROOT/'docs'/'검증스크린샷'; L=S/'2026-08-27T05-19-10-578Z'; A=S/'Codex_T04_독립감사_2026-08-27T02-03-29-907Z'
OLD=S/'T04감사_2026-08-27T00-52-36-485Z'

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BG=RGBColor(10,15,24); PANEL=RGBColor(21,30,44); WHITE=RGBColor(246,248,252); MUTED=RGBColor(162,178,198); GREEN=RGBColor(69,210,163); ORANGE=RGBColor(255,183,78)

def shell(title,source='T04 제출자료 · 지정 문서 기반'):
    s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=BG
    b=s.shapes.add_textbox(Inches(.55),Inches(.2),Inches(8),Inches(.25)); p=b.text_frame.paragraphs[0]; p.text=source; p.font.size=Pt(9); p.font.bold=True; p.font.color.rgb=GREEN
    b=s.shapes.add_textbox(Inches(.55),Inches(.52),Inches(12),Inches(.62)); p=b.text_frame.paragraphs[0]; p.text=title; p.font.size=Pt(25); p.font.bold=True; p.font.color.rgb=WHITE
    b=s.shapes.add_textbox(Inches(12.25),Inches(7.12),Inches(.55),Inches(.2)); p=b.text_frame.paragraphs[0]; p.text=str(len(prs.slides)); p.font.size=Pt(8); p.font.color.rgb=MUTED
    return s

def textbox(s,text,x,y,w,h,size=17,color=WHITE,bold=False,fill=None):
    b=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    if fill: b.fill.solid(); b.fill.fore_color.rgb=fill
    tf=b.text_frame; tf.word_wrap=True; tf.margin_left=Inches(.14); tf.margin_right=Inches(.14); tf.margin_top=Inches(.1)
    p=tf.paragraphs[0]; p.text=text; p.font.size=Pt(size); p.font.color.rgb=color; p.font.bold=bold

def text(title,items,source):
    s=shell(title,source); textbox(s,'\n'.join('• '+i for i in items),.75,1.35,11.85,5.55,18,WHITE,False,PANEL)

def picture(s,path,x=.55,y=1.2,w=12.23,h=5.65):
    with Image.open(path) as im: iw,ih=im.size
    k=min(w/iw,h/ih); pw,ph=iw*k,ih*k
    s.shapes.add_picture(str(path),Inches(x+(w-pw)/2),Inches(y+(h-ph)/2),Inches(pw),Inches(ph))

def evidence(title,path,caption,source):
    s=shell(title,source); picture(s,path); textbox(s,caption,.65,6.85,11.8,.23,10,MUTED)

# 1
s=shell('오늘의 진짜 정보판 — 데이터가 안 올 때','T04 · 최종 제출자료 v2')
textbox(s,'실제 데이터의 값·맥락·실패 상태·일별 변화를\n정직하게 보여주는 공개 대시보드',.75,1.45,7.6,1.35,24,WHITE,True)
textbox(s,'공개 주소\nhttps://realtime-info-board-t04.vercel.app/\n\n소스\ngithub.com/stulss/realtime-info-board-t04',.75,3.3,7.8,2.2,16,WHITE,False,PANEL)
textbox(s,'기획\nC01–C28\nAI 3줄\n검증 체크리스트\n트러블슈팅\n검증 안내서\n작업내역',9.35,1.4,2.8,4.4,17,WHITE,True,PANEL)

text('사용한 문서와 증거 원칙',['지정된 7개 문서만 본문 근거로 사용','검증 주장 바로 뒤에 관련 스크린샷 배치','합성 장애 fixture와 실제 날짜 기록을 분리','완료와 보류를 같은 표에서 명확히 구분'],'기획 · 요구사항 · AI3줄 · 체크리스트 · 트러블슈팅 · 안내서 · 작업내역')
evidence('기획 · 프로젝트 개요',L/'02_실시간_대시보드.png','실시간 관심 데이터를 투명하게 보여주는 개인용 대시보드','docs/01_기획.md')
text('기획 · 기술 스택',['Next.js App Router + React + TypeScript strict','Tailwind CSS + TanStack Query','Vitest + Testing Library + Playwright','Vercel 공개 배포'],'docs/01_기획.md')
evidence('기획 · 위젯 5개',L/'01_상태별_카드_검증.png','Lost Ark 공지·거래장, Upbit, 수출입은행 환율, GitHub Status','docs/01_기획.md')
text('기획 · 시스템 구조',['브라우저는 /api/widgets/* Route Handler만 호출','서버 provider adapter가 외부 API를 조회·정규화','공통 계층: timeout·retry·429·503·TTL cache·dedupe','마지막 성공값은 실패 시 stale로 전환'],'docs/01_기획.md')
text('기획 · 구현 API',['Upbit ticker — 20초 TTL','GitHub Status — 5분','수출입은행 oapi — 1시간','Lost Ark notices — 15분','Lost Ark market — 5분'],'docs/01_기획.md')
text('과제 요구사항 · 전체 판정',['완료: C01–C18, C20, C23–C28','보류: C19·C21·C22','보류 이유: recover-d2 asset과 실제 다음 KST 날짜 재실행','합성 D1/D2를 실제 이틀 기록으로 주장하지 않음'],'docs/00_과제_요구사항_매핑.md')
evidence('C01–C10 · 공개 데이터와 맥락',A/'C1_현재값_단위_출처_조회시각.png','공개 주소·BTC/KRW·값·단위·출처·원천 시각·조회 시각·Asia/Seoul','요구사항 매핑 C01–C10')
evidence('C01–C10 · 화면 증거',L/'02_실시간_대시보드.png','실제 조회 화면의 값과 일별 기록 영역','검증 스크린샷')
evidence('C11 · 비밀 없는 호출',A/'C0_T04_독립감사_메타증거.png','소스·배포 파일·네트워크·Git 비밀값 검색 0건','요구사항 매핑 C11')
evidence('C06 · 출처 한 번 클릭',L/'06_출처링크_한번클릭_원자료.png','화면 출처 링크에서 Upbit 원자료 페이지로 이동','요구사항 매핑 C06')
text('C12–C16 · 장애 5종',['timeout — 지연 후 client Abort','401/403 — 인증 실패','429 — 호출 제한','offline — 연결 실패','schema 변경 — 필수 필드 없는 JSON','모든 장애에는 합성 시험값만 사용'],'검증 체크리스트')
evidence('C12 · timeout',L/'장애_timeout.png','느린 응답을 시간 초과 상태로 구분','검증 체크리스트')
evidence('C13 · 401/403 인증 실패',L/'장애_unauthorized.png','외부 원천 인증 거절을 일반 오류와 구분','검증 체크리스트')
evidence('C14 · 429 호출 제한',L/'장애_rate-limited.png','호출 제한 상태와 다시 시도 행동','검증 체크리스트')
evidence('C15 · 오프라인',L/'장애_offline.png','연결 실패를 오프라인으로 구분','검증 체크리스트')
evidence('C16 · 응답 형식 변경',L/'장애_schema-changed.png','필수 필드 없는 응답을 정상 데이터로 사용하지 않음','검증 체크리스트')
evidence('C17–C18 · 정상값이 없을 때',L/'장애_정상값없음_빈상태.png','숫자를 꾸며내지 않고 빈 상태 표시','검증 체크리스트')
evidence('C19 · 다시 시도와 복구',L/'장애_다시시도_복구.png','복구 후 fresh 상태와 정상 데이터로 전환','검증 체크리스트')
evidence('C20–C24 · KST 일별 기록',L/'05_날짜별기록_중복방지_어제비교.png','같은 날짜 중복 방지와 두 날짜 원천·저장·화면 대조','요구사항 매핑 · 검증 체크리스트')
evidence('실제 두 날짜 · 화면 증거',L/'05_날짜별기록_중복방지_어제비교.png','2026-08-25 109,165,000 → 2026-08-26 109,830,000 KRW','검증 체크리스트')
text('변화값 손계산',['109,830,000 - 109,165,000 = +665,000 KRW','665,000 / 109,165,000 × 100 = 0.6091…%','화면 반올림 +0.61%와 일치','원천=저장=계산 입력=화면값'],'요구사항 매핑 C23–C24')
evidence('원천·저장 불일치 방어',L/'07_원자료_저장값_불일치_비교중단.png','불일치 데이터를 정상 비교값처럼 표시하지 않음','검증 체크리스트')
text('C25–C28 · AI 3줄',['개인정보·비밀값 0건','합성 장애와 실제 기록 분리','검증 안내서 4개 제목·3단계 행동','AI: 장애·일별 기록·테스트·증거 문서','직접 판단: 복사본 작업·실제 Upbit 일봉','불채택: 증거 덮어쓰기·미완료 허위 통과'],'docs/AI_3줄.md · 요구사항 매핑')
text('검증 안내 · 어디로 가나요',['공개 대시보드: realtime-info-board-t04.vercel.app','장애 검증: 같은 주소의 /verification','원본 배포본과 T04 복사본 주소를 구분'],'docs/검증안내서.md')
text('검증 안내 · 3단계 행동',['1. 값·단위·출처·조회 시각·실제 날짜 2건과 링크 확인','2. 이전/현재 값·방향·차이·단위와 중복 방지 확인','3. 장애 5종·오래된 데이터·빈 상태·복구 확인'],'docs/검증안내서.md')
text('검증 안내 · 통과와 안 될 때',['통과: 값의 맥락·KST 두 날짜·비교값·서로 다른 장애 문구','값 오류: trade_price·candle_date_time_kst·KRW 대조','중복: localStorage 날짜 고유키와 Asia/Seoul 확인','장애: 장애 유형·선택 버튼·다시 시도 확인'],'docs/검증안내서.md')
text('트러블슈팅 · 빌드와 브라우저',['npm.ps1 차단 → npm.cmd 사용','prerender document 접근 → 서버 경계 검사','상대시각 hydration → 고정 server snapshot + client clock','예상 503만 구분하고 런타임 오류 검사는 유지'],'docs/트러블슈팅.md')
text('트러블슈팅 · 공급자와 배포',['수출입은행 기존 도메인 종료 → 신규 oapi migration','Lost Ark CategoryCode: 0 → 유효 options 카테고리 순회','없는 기본 아이템 → 실제 조회 가능한 파괴강석','Vercel 전환 404 → 실패·성공 세트 모두 보존'],'docs/트러블슈팅.md')
evidence('트러블슈팅 · 증거 덮어쓰기',A/'C0_T04_독립감사_메타증거.png','UTC 실행시각별 폴더 · 기존 대상 존재 시 중단 · 최초 손실 사실도 기록','docs/트러블슈팅.md')
text('작업내역 · 구현 과정',['전략 분석·계획·표준 문서 구축','공통 WidgetData 계약과 5개 provider','timeout/retry/cache/dedupe/stale fallback','로컬 검증과 Vercel 실데이터 검증','원본 복사 후 T04 카드 3–5 보완'],'작업내역_체크리스트.md')
text('작업내역 · 주요 결정',['Route Handler로 키 필요 API를 서버에서만 처리','위젯 5개 확정·근거 없는 서버 혼잡도 제외','공통 오류 계약으로 429·503·network 분리','증거는 UTC 실행별 누적, 기존 파일 수정·삭제 금지','T04는 원본이 아닌 복사본에서만 작업'],'작업내역_체크리스트.md')
text('작업내역 · 최종 검증',['lint 통과','typecheck 통과','T04 확장 Vitest 41개 통과','production build 통과','Playwright E2E 9개 통과','공개 배포와 GitHub 저장소 확인'],'작업내역_체크리스트.md · T04 검증 체크리스트')
text('최종 상태',['완료: 공개 주소·데이터 맥락·비밀값 0건·장애 5종','완료: stale/empty/retry·중복 방지·실제 원천 대조·손계산','완료: AI 3줄·검증 안내·증거 스크린샷·SHA-256','보류: 실제 다음 KST 날짜 앱 재실행','보류: 제공되지 않은 public asset SHA-256'],'요구사항 매핑 · 검증 체크리스트')

assert len(prs.slides)==35
prs.save(OUT); print(OUT)

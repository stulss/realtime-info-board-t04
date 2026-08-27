from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

ROOT=Path(__file__).resolve().parents[1]
OUT=ROOT/'docs'/'T04_과제_제출_보고서_v3_2026-08-27.pptx'
if OUT.exists(): raise SystemExit(f'refusing to overwrite {OUT}')
S=ROOT/'docs'/'검증스크린샷'; L=S/'2026-08-27T05-19-10-578Z'; A=S/'Codex_T04_독립감사_2026-08-27T02-03-29-907Z'

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BG=RGBColor(55,57,53); CARD=RGBColor(22,25,21); CARD2=RGBColor(44,48,40)
INK=RGBColor(238,238,231); MUTED=RGBColor(178,181,166); LIME=RGBColor(198,255,34)
RED=RGBColor(255,105,92); AMBER=RGBColor(255,181,50); BLUE=RGBColor(112,163,255)

def base(title,kicker='T04 · PULSEBOARD'):
    s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=BG
    tag=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(.58),Inches(.22),Inches(2.35),Inches(.34)); tag.fill.solid(); tag.fill.fore_color.rgb=LIME; tag.line.fill.background()
    p=tag.text_frame.paragraphs[0]; p.text=kicker; p.font.size=Pt(8); p.font.bold=True; p.font.color.rgb=CARD
    t=s.shapes.add_textbox(Inches(.58),Inches(.7),Inches(11.9),Inches(.65)); p=t.text_frame.paragraphs[0]; p.text=title; p.font.size=Pt(25); p.font.bold=True; p.font.color.rgb=INK
    n=s.shapes.add_textbox(Inches(12.15),Inches(7.06),Inches(.6),Inches(.25)); p=n.text_frame.paragraphs[0]; p.text=f'{len(prs.slides):02d}'; p.font.size=Pt(8); p.font.color.rgb=MUTED
    return s

def card(s,x,y,w,h,fill=CARD):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h)); sh.fill.solid(); sh.fill.fore_color.rgb=fill; sh.line.color.rgb=RGBColor(75,78,69); return sh

def tx(s,text,x,y,w,h,size=16,color=INK,bold=False):
    b=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=b.text_frame; tf.word_wrap=True; tf.margin_left=Inches(.1); tf.margin_right=Inches(.1)
    p=tf.paragraphs[0]; p.text=text; p.font.size=Pt(size); p.font.color.rgb=color; p.font.bold=bold

def image(s,path,x,y,w,h):
    with Image.open(path) as im: iw,ih=im.size
    k=min(w/iw,h/ih); pw,ph=iw*k,ih*k
    s.shapes.add_picture(str(path),Inches(x+(w-pw)/2),Inches(y+(h-ph)/2),Inches(pw),Inches(ph))

def evidence(title,path,caption,kicker='EVIDENCE · SCREENSHOT'):
    s=base(title,kicker); card(s,.58,1.42,12.15,5.45); image(s,path,.73,1.58,11.85,5.0); tx(s,caption,.78,6.66,11.7,.25,10,MUTED)

def cards(title,left_title,left,right_title,right,kicker='DOCUMENT · SUMMARY'):
    s=base(title,kicker); card(s,.65,1.55,5.92,4.95); card(s,6.77,1.55,5.92,4.95)
    tx(s,left_title,.9,1.82,5.3,.35,13,LIME,True); tx(s,'\n'.join('• '+v for v in left),.9,2.35,5.3,3.75,17,INK)
    tx(s,right_title,7.02,1.82,5.3,.35,13,LIME,True); tx(s,'\n'.join('• '+v for v in right),7.02,2.35,5.3,3.75,17,INK)

def list_slide(title,items,kicker='DOCUMENT · SUMMARY'):
    s=base(title,kicker); card(s,.65,1.5,12.04,5.05); tx(s,'\n'.join('• '+v for v in items),.95,1.88,11.35,4.25,18,INK)

# 1 cover
s=base('오늘의 진짜 정보판 — 데이터가 안 올 때','T04 · 과제 제출 보고서')
tx(s,'실제 데이터가 늦거나 실패해도\n마지막 정상값과 현재 상태를 정직하게 보여주는 정보판',.72,1.55,7.3,1.6,24,INK,True)
card(s,.72,3.55,7.45,2.05); tx(s,'결과물  realtime-info-board-t04.vercel.app\n소스     github.com/stulss/realtime-info-board-t04\n데이터   Upbit BTC/KRW · Asia/Seoul',1.02,3.94,6.85,1.3,15,INK)
card(s,9.0,1.5,3.1,4.8,CARD2); tx(s,'5',9.5,1.95,1.2,.7,36,LIME,True); tx(s,'과제 카드',10.25,2.15,1.45,.35,13,MUTED,True); tx(s,'5종\n장애 재현',9.5,3.05,2.0,1.0,23,INK,True); tx(s,'2일\n실제 원천 기록',9.5,4.45,2.0,1.0,23,INK,True)

# 2 overview
cards('OVERVIEW · 프로젝트 개요','문제',['현재값만으로 출처·시각을 알 수 없음','실패가 같은 오류로 보이면 대응 불가','기록과 계산이 원천과 같은지 검증 어려움'],'해결',['값·단위·출처·두 시각·KST','장애 종류·stale·empty·retry','일별 기록과 어제 대비 손계산'],'기획 · 요구사항')

# 3 guide
s=base('USER GUIDE · 30초 확인 방법','검증안내서')
for i,(head,body) in enumerate([('01 · 메인','값·단위·출처·두 시각·KST'),('02 · 기록','실제 2행·변화값·출처 클릭'),('03 · 장애','5종·빈 상태·다시 시도')]):
    x=.72+i*4.08; card(s,x,1.65,3.78,4.55); tx(s,head,x+.25,2.0,3.2,.4,14,LIME,True); tx(s,body,x+.25,2.75,3.25,1.4,19,INK,True); tx(s,'설치·로그인 없이\n공개 주소에서 확인',x+.25,5.1,3.1,.7,12,MUTED)

# 4-7 cards 1-2
cards('CARD 1 · 매일 궁금한 값 하나','사용 데이터',['Upbit 공개 BTC/KRW 완료 일봉','trade_price · candle_date_time_kst','기준 시간대 Asia/Seoul'],'화면 맥락',['현재값·KRW 단위','원천 URL·원천 관측 시각','조회 시각·일별 기록'],'과제 요구사항 C03–C10')
evidence('CARD 1 · 값의 맥락 증거',A/'C1_현재값_단위_출처_조회시각.png','현재값·단위·출처·원천 시각·조회 시각을 한 화면에서 확인','CARD 1 · EVIDENCE')
cards('CARD 2 · 비밀 없는 호출','호출 경로',['브라우저 → /api/widgets/*','서버 provider → 외부 API','WidgetData<T> 정규화'],'검사 결과',['소스·배포 파일 비밀값 0건','네트워크·Git 비밀값 0건','브라우저 provider 직접 요청 0건'],'과제 요구사항 C11')
evidence('CARD 2 · 출처 한 번 클릭',L/'06_출처링크_한번클릭_원자료.png','Upbit 출처 링크를 한 번 눌러 원자료 페이지로 이동','CARD 2 · EVIDENCE')

# 8-15 card3
cards('CARD 3 · 다섯 가지 실패','합성 재현',['timeout · 401/403 · 429','offline · schema 변경','실제 원천/날짜 기록과 분리'],'공통 행동',['오류 종류별 문구','마지막 정상값은 stale','정상값 없으면 empty · 다시 시도'],'검증 체크리스트 C12–C19')
evidence('CARD 3 · timeout',L/'장애_timeout.png','느린 응답을 별도 시간 초과 상태로 표시','C12 · EVIDENCE')
evidence('CARD 3 · 인증 실패',L/'장애_unauthorized.png','HTTP 401/403을 인증 실패로 구분','C13 · EVIDENCE')
evidence('CARD 3 · 호출 제한',L/'장애_rate-limited.png','HTTP 429를 호출 제한으로 구분','C14 · EVIDENCE')
evidence('CARD 3 · 오프라인',L/'장애_offline.png','연결 실패를 별도 오프라인 상태로 표시','C15 · EVIDENCE')
evidence('CARD 3 · 응답 형식 변경',L/'장애_schema-changed.png','필수 필드 없는 JSON을 정상값으로 사용하지 않음','C16 · EVIDENCE')
evidence('CARD 3 · 정상값이 없을 때',L/'장애_정상값없음_빈상태.png','값을 꾸며내지 않고 표시할 데이터 없음으로 안내','C17–C18 · EVIDENCE')
evidence('CARD 3 · 다시 시도 복구',L/'장애_다시시도_복구.png','다시 시도 후 fresh 상태와 정상 데이터로 복구','C19 · EVIDENCE')

# 16-19 cards4-5
cards('CARD 4 · 하루 한 줄','고유키',['Asia/Seoul 날짜','데이터 종류','upbit:KRW-BTC:YYYY-MM-DD'],'갱신 규칙',['같은 날짜는 한 행으로 병합','다른 날짜는 새 행','합성 기록과 실제 기록 분리'],'과제 요구사항 C20–C22')
evidence('CARD 4 · 날짜 중복 방지 증거',L/'05_날짜별기록_중복방지_어제비교.png','같은 날짜 동기화 후에도 날짜별 고유 행 유지','CARD 4 · EVIDENCE')
cards('CARD 5 · 실제 이틀과 어제 대비','실제 원천 기록',['2026-08-25 · 109,165,000 KRW','2026-08-26 · 109,830,000 KRW','원천=저장=계산=화면'],'손계산',['차이 +665,000 KRW','변화율 0.6091…%','화면 반올림 +0.61%'],'과제 요구사항 C23–C24')
evidence('CARD 5 · 비교와 방어 증거',L/'07_원자료_저장값_불일치_비교중단.png','원천과 저장값이 다르면 비교 계산을 중단','CARD 5 · EVIDENCE')

# 20-21 tech and architecture
cards('TECH DECISION · 왜 Route Handler인가','브라우저 직접 호출의 위험',['비밀키·공급자 주소 노출','provider별 오류 처리 중복','응답 형식 변경이 UI까지 전파'],'서버 경계의 이점',['비밀값 서버 전용','공통 오류·스키마 검증','cache·dedupe·stale 재사용'],'docs/01_기획.md')
cards('ARCHITECTURE · 시스템과 위젯','기술 스택',['Next.js · React · TypeScript strict','TanStack Query · Tailwind','Vitest · Playwright · Vercel'],'5개 위젯',['Lost Ark 공지·거래장','Upbit · 수출입은행 환율','GitHub Status'],'docs/01_기획.md')

# 22-24 troubleshooting
cards('TROUBLESHOOTING 1 · SSR과 브라우저','문제',['npm.ps1 실행 정책 차단','prerender에서 document 접근','상대시각 hydration 불일치'],'해결',['npm.cmd 사용','typeof document 서버 경계','고정 snapshot + client clock'],'docs/트러블슈팅.md')
cards('TROUBLESHOOTING 2 · 공급자와 배포','문제',['수출입은행 기존 도메인 종료','Lost Ark CategoryCode·기본 아이템','Vercel 전환 중 일시 404'],'해결',['신규 oapi migration','options 순회·파괴강석','실패/성공 증거 모두 보존'],'docs/트러블슈팅.md')
evidence('TROUBLESHOOTING 3 · 증거 덮어쓰기',A/'C0_T04_독립감사_메타증거.png','UTC 실행별 새 폴더 · 기존 대상 존재 시 중단 · 최초 손실 사실 기록','docs/트러블슈팅.md')

# 25 AI
list_slide('AI 3줄',['AI에게 맡긴 일 — 기존 코드 보존, 장애 5종, 일별 기록, 비교, 테스트, 증거 문서','직접 판단한 일 — 복사본에서만 작업, 실제 Upbit 일봉 사용','AI 말을 안 들은 일 — 기존 증거 덮어쓰기, 미완료 항목 허위 통과'],'docs/AI_3줄.md')

# 26-27 guide
cards('VERIFICATION GUIDE 1 · 어디로 / 무엇을','어디로 가나요',['realtime-info-board-t04.vercel.app','같은 주소의 /verification','원본 프로젝트 주소와 구분'],'무엇을 하나요',['값의 맥락 확인','두 날짜·변화·출처 확인','장애 5종·빈 상태·복구'],'docs/검증안내서.md')
cards('VERIFICATION GUIDE 2 · 통과 / 안 될 때','통과 모습',['값·KRW·출처·두 시각','KST 실제 두 행·변화값','장애별 문구·stale·empty'],'안 될 때',['Upbit 원천 필드 대조','localStorage 날짜 키 확인','장애 유형·다시 시도 확인'],'docs/검증안내서.md')

# 28-29 checklist
cards('CHECKLIST · C01–C16','C01–C11',['공개 주소·실제 원천','값·단위·출처·두 시각·KST','원천/저장/화면 일치·비밀값 0건'],'C12–C16',['timeout · 401/403 · 429','offline · schema 변경','합성 시험값만 사용'],'요구사항 매핑')
cards('CHECKLIST · C17–C28','완료',['C17–C18 stale/empty · C20 중복','C23–C24 대조·손계산','C25–C28 개인정보·안내·AI 3줄'],'보류',['C19 공개 recover-d2 asset','C21 다음 실제 KST 날짜','C22 서로 다른 날 앱 직접 실행'],'요구사항 매핑')

# 30-32 work/final
list_slide('WORK LOG · 구현 과정',['전략 분석 → 계획과 표준 문서','WidgetData 계약 → 5개 provider','cache·error·dedupe·stale fallback','lint·typecheck·unit·build·E2E','Vercel 배포 → T04 복사본 카드 3–5 보완'],'작업내역_체크리스트.md')
cards('FINAL STATUS','완료',['공개 주소·데이터 맥락·비밀값 0건','장애 5종·stale/empty/retry','KST 중복·실제 원천 대조·손계산'],'보류',['실제 다음 KST 날짜 앱 재실행','제공되지 않은 public asset SHA-256','증거를 조작하지 않고 보류 표시'],'검증 체크리스트')
s=base('제출 주소','T04 · FINAL LINKS'); card(s,.75,1.55,11.8,4.85); tx(s,'결과물',1.05,2.0,1.4,.35,13,LIME,True); tx(s,'https://realtime-info-board-t04.vercel.app/',2.55,1.95,8.9,.45,18,INK)
tx(s,'장애 검증',1.05,3.05,1.4,.35,13,LIME,True); tx(s,'https://realtime-info-board-t04.vercel.app/verification',2.55,3.0,9.0,.45,18,INK)
tx(s,'소스',1.05,4.1,1.4,.35,13,LIME,True); tx(s,'https://github.com/stulss/realtime-info-board-t04',2.55,4.05,9.0,.45,18,INK)
tx(s,'원고',1.05,5.15,1.4,.35,13,LIME,True); tx(s,'docs/T04_제출자료_원고_v3.md',2.55,5.1,9.0,.45,18,INK)

assert len(prs.slides)==32
prs.save(OUT); print(OUT)

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "T04_발표자료_2026-08-27.pptx"
if OUT.exists():
    raise SystemExit(f"refusing to overwrite {OUT}")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BG = RGBColor(14, 20, 32)
WHITE = RGBColor(245, 247, 250)
MUTED = RGBColor(170, 182, 199)
ACCENT = RGBColor(90, 210, 170)

def slide(title, subtitle=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    box = s.shapes.add_textbox(Inches(.7), Inches(.45), Inches(12), Inches(.8))
    tf = box.text_frame; tf.clear(); p = tf.paragraphs[0]
    p.text = title; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = WHITE
    if subtitle:
        b = s.shapes.add_textbox(Inches(.72), Inches(1.2), Inches(11.8), Inches(.45))
        p = b.text_frame.paragraphs[0]; p.text = subtitle; p.font.size = Pt(13); p.font.color.rgb = MUTED
    return s

def bullets(s, items, y=1.9, size=20):
    b = s.shapes.add_textbox(Inches(1), Inches(y), Inches(11.2), Inches(4.8))
    tf = b.text_frame; tf.word_wrap = True; tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(size); p.font.color.rgb = WHITE; p.space_after = Pt(14)

def image_slide(title, image, caption):
    s = slide(title, caption)
    path = ROOT / "docs" / "검증스크린샷" / "2026-08-27T05-19-10-578Z" / image
    if path.exists():
        s.shapes.add_picture(str(path), Inches(.85), Inches(1.65), width=Inches(11.7), height=Inches(5.25))
    else:
        bullets(s, [f"증거 이미지 없음: {image}"], y=2.2)
    return s

s = slide("오늘의 진짜 정보판", "T04 발표자료 · 데이터가 안 올 때")
bullets(s, ["실제로 변하는 공개 데이터를 매일 기록합니다.", "값이 늦거나 실패해도 마지막 정상값과 현재 상태를 정직하게 보여줍니다.", "공개 결과: https://realtime-info-board-t04.vercel.app/"], y=2.0, size=22)

s = slide("무엇을 만들었나", "Upbit BTC/KRW 일봉 기반 개인 정보판")
bullets(s, ["값 · 단위(KRW) · 출처 링크 · 원천 관측 시각", "조회 시각 · Asia/Seoul 기준 시간대 · 일별 기록", "어제 대비 차이·방향·변화율을 한 화면에서 확인"])

s = slide("검증 기준", "변경된 T04 C01–C28을 기준으로 판정")
bullets(s, ["카드 1: 값의 맥락과 원자료 대조", "카드 2: 브라우저·배포 파일·네트워크·Git 비밀값 0건", "카드 3: timeout·401/403·429·offline·schema 변경", "카드 4–5: KST 일별 중복 방지와 어제 대비 계산"])

image_slide("정상 화면과 일별 비교", "05_날짜별기록_중복방지_어제비교.png", "원천값·저장값·계산 입력·화면값 대조")
image_slide("출처 링크 증거", "06_출처링크_한번클릭_원자료.png", "한 번 클릭해 원자료 페이지로 이동")
image_slide("응답 형식 변경", "장애_schema-changed.png", "실패 종류를 별도 상태로 표시")
image_slide("정상값 보존과 복구", "장애_다시시도_복구.png", "다시 시도 후 정상 상태로 복구")

s = slide("실제 원천 기록 2건", "합성 fixture는 실제 기록을 대신하지 않음")
bullets(s, ["2026-08-25: 109,165,000 KRW", "2026-08-26: 109,830,000 KRW", "차이: +665,000 KRW · 변화율: +0.61%", "두 날짜의 원천·저장·계산·화면값이 일치"])

s = slide("안전한 호출 경로", "서버 Route Handler → 검증된 응답 → 화면")
bullets(s, ["브라우저에서 외부 provider를 직접 호출하지 않습니다.", "응답 스키마와 HTTP 상태를 경계에서 검증합니다.", "실패 시 마지막 정상값을 보존하고 stale/empty를 구분합니다.", "비밀값 검색 결과: 0건"])

s = slide("검증 결과", "로컬 자동 검증")
bullets(s, ["lint 통과", "typecheck 통과", "Vitest 41개 통과", "production build 통과", "Playwright E2E 9개 통과", "증거 폴더는 실행 timestamp별로 생성되며 덮어쓰기 방지"])

s = slide("남은 확인과 정직한 판정")
bullets(s, ["공개 배포 URL은 완료: https://realtime-info-board-t04.vercel.app/", "다음 실제 KST 날짜의 앱 재실행(C19/C21/C22)은 추가 확인 항목", "제공되지 않은 공개 asset 패키지 SHA-256도 보류", "합성 D1/D2를 실제 날짜 기록처럼 사용하지 않음"])

s = slide("마무리", "소스·문서·스크린샷 증거를 함께 제출")
bullets(s, ["소스: https://github.com/stulss/realtime-info-board-t04", "전체 원고: docs/T04_프로젝트_전체기록.md", "검증 안내: docs/T04_검증_체크리스트.md", "질문 감사합니다."])

prs.save(OUT)
print(OUT)

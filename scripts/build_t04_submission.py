from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "T04_제출자료_2026-08-27.pptx"
if OUT.exists():
    raise SystemExit(f"refusing to overwrite {OUT}")

LATEST = ROOT / "docs" / "검증스크린샷" / "2026-08-27T05-19-10-578Z"
AUDIT = ROOT / "docs" / "검증스크린샷" / "Codex_T04_독립감사_2026-08-27T02-03-29-907Z"

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BG=RGBColor(9,14,23); PANEL=RGBColor(20,28,42); WHITE=RGBColor(245,248,252)
MUTED=RGBColor(165,180,199); ACCENT=RGBColor(67,211,163); WARN=RGBColor(255,185,80)

def page(title, eyebrow="T04 · 오늘의 진짜 정보판"):
    s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=BG
    b=s.shapes.add_textbox(Inches(.55),Inches(.22),Inches(8),Inches(.28)); p=b.text_frame.paragraphs[0]
    p.text=eyebrow; p.font.size=Pt(9); p.font.bold=True; p.font.color.rgb=ACCENT
    b=s.shapes.add_textbox(Inches(.55),Inches(.52),Inches(12),Inches(.65)); p=b.text_frame.paragraphs[0]
    p.text=title; p.font.size=Pt(25); p.font.bold=True; p.font.color.rgb=WHITE
    b=s.shapes.add_textbox(Inches(12.25),Inches(7.12),Inches(.55),Inches(.2)); p=b.text_frame.paragraphs[0]
    p.text=str(len(prs.slides)); p.font.size=Pt(8); p.font.color.rgb=MUTED; p.alignment=PP_ALIGN.RIGHT
    return s

def box(s,text,x,y,w,h,size=16,color=WHITE,bold=False,fill=None):
    b=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    if fill: b.fill.solid(); b.fill.fore_color.rgb=fill
    tf=b.text_frame; tf.word_wrap=True; tf.margin_left=Inches(.14); tf.margin_right=Inches(.14)
    p=tf.paragraphs[0]; p.text=text; p.font.size=Pt(size); p.font.color.rgb=color; p.font.bold=bold
    return b

def picture(s,path,x=.55,y=1.25,w=12.23,h=5.55):
    with Image.open(path) as im: iw,ih=im.size
    scale=min(w/iw,h/ih); pw,ph=iw*scale,ih*scale
    s.shapes.add_picture(str(path),Inches(x+(w-pw)/2),Inches(y+(h-ph)/2),Inches(pw),Inches(ph))

def evidence(title,path,caption):
    s=page(title); picture(s,path); box(s,caption,.65,6.82,11.9,.25,10,MUTED)

def bullets(title,items,subtitle=""):
    s=page(title)
    if subtitle: box(s,subtitle,.7,1.25,11.8,.45,13,MUTED)
    text="\n".join(f"• {v}" for v in items); box(s,text,.85,1.85,11.6,4.9,19,WHITE,False,PANEL)

# 1
s=page("데이터가 안 올 때도 정직한 정보판")
box(s,"Upbit BTC/KRW 일봉을 매일 기록하고\n실패에도 마지막 정상값과 현재 상태를 설명합니다.",.8,1.45,7.4,1.45,25,WHITE,True)
box(s,"공개 결과\nhttps://realtime-info-board-t04.vercel.app/\n\n소스\ngithub.com/stulss/realtime-info-board-t04",.8,3.35,7.6,2.1,16,WHITE,False,PANEL)
box(s,"실제 공개 데이터\n합성 장애 5종\nKST 일별 기록\n스크린샷 SHA-256",9.3,1.6,2.8,3.75,18,WHITE,True,PANEL)

# 2-4
evidence("프로젝트 개요",LATEST/"02_실시간_대시보드.png","값·출처·시각·실패 상태·어제 대비 변화를 한 화면에 구성")
bullets("기능 안내",["/ — 실시간 카드, 일별 기록, 비교값, 출처 링크","/verification — 다섯 합성 장애, stale/empty, 다시 시도","Upbit 원자료 — 실제 일봉 값과 KST 관측 시각"])
bullets("데이터 호출 구조",["브라우저 → Next.js Route Handler","provider adapter → Upbit 공개 API","응답 스키마 검증 → 화면·일별 기록","브라우저는 외부 provider를 직접 호출하지 않음"],"비밀값을 브라우저·배포 파일·네트워크 응답에 남기지 않는 서버 경계")

# 5-9
evidence("CARD 1 · 값·단위·출처·두 시각",AUDIT/"C1_현재값_단위_출처_조회시각.png","현재값·KRW·Upbit 출처·원천 관측 시각·조회 시각을 같은 화면에서 확인")
evidence("공개 대시보드 전체",LATEST/"02_실시간_대시보드.png","실제 공개 원천과 T04 일별 기록 영역")
evidence("상태별 카드",LATEST/"01_상태별_카드_검증.png","fresh · stale · provider error · empty를 구분")
evidence("CARD 2 · 출처를 한 번 클릭",LATEST/"06_출처링크_한번클릭_원자료.png","Upbit 공개 원자료 페이지가 새 탭에서 열림")
evidence("비밀 없는 호출",AUDIT/"C0_T04_독립감사_메타증거.png","브라우저 외부 요청·소스·배포 파일·Git 비밀값 검사 0건")

# 10-17
evidence("CARD 3 · 다섯 실패 검증",LATEST/"장애_timeout.png","모든 실패는 합성 시험값이며 실제 날짜 기록을 대체하지 않음")
evidence("실패 1 · timeout",LATEST/"장애_timeout.png","느린 응답을 시간 초과로 분리하고 다시 시도 제공")
evidence("실패 2 · 401/403 인증 거절",LATEST/"장애_unauthorized.png","외부 원천 인증 실패를 일반 오류와 구분")
evidence("실패 3 · 429 호출 제한",LATEST/"장애_rate-limited.png","호출 제한 상태와 다음 행동을 별도 표시")
evidence("실패 4 · 오프라인",LATEST/"장애_offline.png","연결 실패를 별도 오프라인 상태로 표시")
evidence("실패 5 · 응답 형식 변경",LATEST/"장애_schema-changed.png","필수 필드가 없는 JSON을 정상 데이터로 사용하지 않음")
evidence("정상값이 없으면 빈 상태",LATEST/"장애_정상값없음_빈상태.png","값을 꾸며내지 않고 표시할 데이터 없음으로 안내")
evidence("다시 시도 후 정상 복구",LATEST/"장애_다시시도_복구.png","복구 후 fresh 상태와 정상 데이터로 돌아옴")

# 18-21
evidence("CARD 4 · 하루 한 줄",LATEST/"05_날짜별기록_중복방지_어제비교.png","Asia/Seoul 날짜+데이터 종류 고유키로 같은 날 중복 방지")
evidence("CARD 5 · 실제 두 날짜 대조",LATEST/"05_날짜별기록_중복방지_어제비교.png","2026-08-25 109,165,000 → 2026-08-26 109,830,000 KRW")
bullets("어제 대비 손계산",["109,830,000 - 109,165,000 = +665,000 KRW","(665,000 / 109,165,000) × 100 = 0.6091…%","화면 반올림 결과 +0.61%와 일치","방향·차이·단위·변화율을 같은 계산 규칙으로 재검산"])
evidence("원자료와 저장값이 다르면 중단",LATEST/"07_원자료_저장값_불일치_비교중단.png","불일치 데이터를 정상 비교값처럼 표시하지 않는 방어 동작")

# 22-25
bullets("자동 검증",["ESLint 통과","TypeScript typecheck 통과","Vitest 41개 통과","production build 통과","Playwright E2E 9개 통과","스크린샷별 SHA-256 기록"])
evidence("증거 보관 · 덮어쓰기 방지",AUDIT/"C0_T04_독립감사_메타증거.png","실행마다 UTC timestamp 폴더 생성 · 기존 대상이 있으면 중단")
bullets("30초 검증 안내",["1. 공개 주소에서 값·단위·출처·두 시각·KST 확인","2. 일별 두 행과 변화값 확인 후 출처 링크 한 번 클릭","3. /verification에서 장애 5종·빈 상태·복구 확인","통과: 현재/오래된/빈 상태와 실패 문구가 구분되고 계산이 일치"])
s=page("최종 상태와 AI 판단")
box(s,"완료",.75,1.35,1.2,.35,13,ACCENT,True); box(s,"공개 주소 · 데이터 맥락 · 비밀값 0건\n장애 5종 · stale/empty/retry · 중복 방지\n실제 원천 2건 대조 · 문서와 증거",.75,1.75,5.75,2.0,17,WHITE,False,PANEL)
box(s,"보류",6.85,1.35,1.2,.35,13,WARN,True); box(s,"실제 다음 KST 날짜 앱 재실행\n제공되지 않은 public asset SHA-256\n합성값은 실제 기록으로 주장하지 않음",6.85,1.75,5.7,2.0,17,WHITE,False,PANEL)
box(s,"AI에게 맡긴 일",.75,4.25,2,.3,11,ACCENT,True); box(s,"오류 분기·일별 기록·자동 검증·증거 문서화",2.65,4.22,9.7,.4,14,WHITE)
box(s,"직접 판단한 일",.75,4.85,2,.3,11,ACCENT,True); box(s,"공개 원천 선택, KST·단위·계산 규칙 대조",2.65,4.82,9.7,.4,14,WHITE)
box(s,"따르지 않은 제안",.75,5.45,2,.3,11,ACCENT,True); box(s,"합성 값을 실제 날짜 기록처럼 저장하지 않음",2.65,5.42,9.7,.4,14,WHITE)
box(s,"전체 원고: docs/T04_제출자료_원고.md",.8,6.5,11.8,.3,11,MUTED)

# 26-32 · 저장소의 나머지 프로젝트 문서 반영
bullets("기획 배경과 사용자",["대상: 매일 BTC/KRW 움직임을 빠르게 확인하려는 사용자","문제: 현재값만으로는 출처·시각·실패 여부를 판단하기 어려움","해결: 데이터 맥락·실패 상태·일별 변화·검증 근거를 한 화면에 결합","기존 대시보드 구조를 보존하고 T04 요구 기능만 확장"],"README.md · 계획.md · docs/01_기획.md")
bullets("기술 선택과 구조",["Next.js + TypeScript + React + TanStack Query","Route Handler로 비밀값과 외부 호출 경계 관리","provider adapter로 응답 변환·검증 책임 분리","공통 cache·dedupe·timeout·retry·stale fallback 재사용","localStorage에는 공개 일별 기록만 저장"],"계획.md · Codex_작업지시서.md · 작업내역_체크리스트.md")
evidence("배포와 운영",LATEST/"02_실시간_대시보드.png","공개 주소 · 서버 전용 환경변수 · 새 시크릿 창 무로그인 접근")
evidence("트러블슈팅과 복구",LATEST/"장애_다시시도_복구.png","HTTP 상태 분리 · offline/timeout 분리 · schema 경계 검증 · 실제 timeout은 보류 기록")
evidence("작업 원칙 · 모바일 검증",ROOT/"docs"/"검증스크린샷"/"T04감사_2026-08-27T00-52-36-485Z"/"X1-모바일_레이아웃.png","원본 복사본 작업 · 기존 코드/테스트 보존 · 증거 덮어쓰기 금지")
evidence("문서·증거 색인 · 다크모드",ROOT/"docs"/"검증스크린샷"/"T04감사_2026-08-27T00-52-36-485Z"/"X2-다크모드.png","기획·배포·체크리스트·검증안내·AI 3줄·트러블슈팅·스크린샷·SHA-256")
bullets("포트폴리오 요약과 결론",["정상 데이터만 꾸미는 화면이 아니라 실패를 종류별로 설명하는 정보 제품","마지막 정상값 보존과 실제 기록 변화 재검산","자동 테스트와 덮어쓰기 불가 스크린샷 증거","완료와 보류를 분리하고 합성값을 실제 기록으로 주장하지 않음"],"docs/포트폴리오_추가용_소개글.md · docs/AI_3줄.md")

assert len(prs.slides)==32
prs.save(OUT); print(OUT)
